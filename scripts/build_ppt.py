#!/usr/bin/env python
"""
Builds paper/MemoryLifeBench_final_defense.pptx: a full, presentation-ready
PowerPoint deck (native, editable slides -- text, tables, and real
PowerPoint charts, not images) covering the whole project for a final
post-graduation-project defense in front of a professor.

Categorical colors follow the validated default palette from the dataviz
skill (fixed hue order, colorblind-safe adjacent pairs): blue, orange,
aqua, yellow, magenta, green, violet, red. Color is assigned by entity
identity (which policy/method) and held constant across every chart in
the deck, never reassigned by rank.

    python scripts/build_ppt.py
"""
import re

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_TICK_MARK

# ---------------------------------------------------------------- palette --

BLUE = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
AQUA = RGBColor(0x1B, 0xAF, 0x7A)
YELLOW = RGBColor(0xED, 0xA1, 0x00)
MAGENTA = RGBColor(0xE8, 0x7B, 0xA4)
GREEN = RGBColor(0x00, 0x83, 0x00)
VIOLET = RGBColor(0x4A, 0x3A, 0xA7)
RED = RGBColor(0xE3, 0x49, 0x48)
CAT8 = [BLUE, ORANGE, AQUA, YELLOW, MAGENTA, GREEN, VIOLET, RED]

INK = RGBColor(0x0B, 0x0B, 0x0B)
INK_SECONDARY = RGBColor(0x52, 0x51, 0x4E)
INK_MUTED = RGBColor(0x89, 0x87, 0x81)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
PAGE = RGBColor(0xF9, 0xF9, 0xF7)
GRIDLINE = RGBColor(0xE1, 0xE0, 0xD9)
GOOD = RGBColor(0x00, 0x63, 0x00)
CRITICAL = RGBColor(0xD0, 0x3B, 0x3B)

# Fixed identity -> color assignment, held constant across the whole deck.
POLICY_COLOR = {
    "no_forget": BLUE,
    "fifo": ORANGE,
    "lru": AQUA,
    "ours": YELLOW,
    "ours_utility": MAGENTA,
    "mem0": GREEN,
    "oracle": VIOLET,
    "full_context": RED,
}
BASELINE_COLOR = {
    "Our model": BLUE,
    "Bucket classifier": ORANGE,
    "GPT-4o": AQUA,
    "Qwen2.5-7B": YELLOW,
    "Gemini 2.5 Pro": MAGENTA,
    "Recency heuristic": GREEN,
}

FONT = "Calibri"
SW, SH = Inches(13.333), Inches(7.5)

INLINE_RE = re.compile(r"(\*\*.+?\*\*)")


def rgb_to_hex_dist(c):
    return c


# --------------------------------------------------------------- helpers --

def new_slide(prs, bg=PAGE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0.75)):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w
    return shp


def add_textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def add_inline_runs(paragraph, text, size, color, italic=False, font=FONT):
    for chunk in INLINE_RE.split(text):
        if not chunk:
            continue
        run = paragraph.add_run()
        if chunk.startswith("**") and chunk.endswith("**"):
            run.text = chunk[2:-2]
            run.font.bold = True
        else:
            run.text = chunk
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = font
        run.font.italic = italic


def kicker_and_title(slide, kicker, title, title_size=28):
    """Small accent kicker label + big slide title, with an accent rule."""
    _, tf = add_textbox(slide, Inches(0.55), Inches(0.32), Inches(11.5), Inches(0.3))
    p = tf.paragraphs[0]
    add_inline_runs(p, kicker.upper(), 12, BLUE)
    p.runs[0].font.bold = True
    p.runs[0].font.name = FONT
    for r in p.runs:
        r.font.spacing = None

    _, tf2 = add_textbox(slide, Inches(0.55), Inches(0.62), Inches(12.2), Inches(0.75))
    p2 = tf2.paragraphs[0]
    add_inline_runs(p2, title, title_size, INK)
    for r in p2.runs:
        r.font.bold = True

    add_rect(slide, Inches(0.55), Inches(1.32), Inches(1.1), Pt(3), fill=BLUE)


def add_bullets(slide, items, x, y, w, h, size=16, color=INK, space_after=8,
                 anchor=MSO_ANCHOR.TOP):
    """items: list of dicts {text, level(0/1), bold(bool), color(optional)}"""
    _, tf = add_textbox(slide, x, y, w, h, anchor=anchor)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        level = item.get("level", 0)
        p.level = level
        bullet = "" if item.get("no_bullet") else ("–  " if level else "▪  ")
        txt = bullet + item["text"]
        c = item.get("color", color)
        add_inline_runs(p, txt, item.get("size", size), c)
        p.space_after = Pt(space_after)
        if item.get("bold"):
            for r in p.runs:
                r.font.bold = True
    return tf


def add_footer(slide, idx, total, section=""):
    add_rect(slide, 0, SH - Pt(3), SW, Pt(3), fill=BLUE)
    _, tf = add_textbox(slide, Inches(0.55), SH - Inches(0.38), Inches(6), Inches(0.3))
    p = tf.paragraphs[0]
    add_inline_runs(p, "MemoryLifeBench  ·  Final Project Defense", 9, INK_MUTED)
    _, tf2 = add_textbox(slide, SW - Inches(1.3), SH - Inches(0.38), Inches(0.8), Inches(0.3))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    add_inline_runs(p2, f"{idx} / {total}", 9, INK_MUTED)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_table(slide, x, y, w, h, headers, rows, col_widths=None, header_fill=BLUE,
              font_size=13, header_color=WHITE, zebra=True, highlight_rows=None,
              highlight_color=RGBColor(0xFF, 0xF5, 0xD6)):
    highlight_rows = highlight_rows or set()
    n_rows, n_cols = len(rows) + 1, len(headers)
    gt = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    table = gt.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Emu(int(w * cw / total))
    for c, htxt in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        add_inline_runs(p, htxt, font_size, header_color)
        for r in p.runs:
            r.font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.08)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)
    for r_i, row in enumerate(rows):
        for c_i, val in enumerate(row):
            cell = table.cell(r_i + 1, c_i)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            add_inline_runs(p, str(val), font_size, INK)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            if r_i in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = highlight_color
            elif zebra and r_i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = SURFACE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return table


def style_chart_axes(chart, y_min=None, y_max=None, y_fmt="0.00", show_val_axis=True):
    cat_ax = chart.category_axis
    cat_ax.tick_labels.font.size = Pt(12)
    cat_ax.tick_labels.font.color.rgb = INK_SECONDARY
    cat_ax.format.line.color.rgb = GRIDLINE
    cat_ax.major_tick_mark = XL_TICK_MARK.NONE
    val_ax = chart.value_axis
    val_ax.visible = show_val_axis
    if show_val_axis:
        val_ax.tick_labels.font.size = Pt(10)
        val_ax.tick_labels.font.color.rgb = INK_MUTED
        val_ax.tick_labels.number_format = y_fmt
        val_ax.tick_labels.number_format_is_linked = False
        val_ax.major_gridlines.format.line.color.rgb = GRIDLINE
        val_ax.major_gridlines.format.line.width = Pt(0.5)
        val_ax.format.line.fill.background()
        if y_min is not None:
            val_ax.minimum_scale = y_min
        if y_max is not None:
            val_ax.maximum_scale = y_max


def add_single_series_bar(slide, x, y, w, h, categories, values, colors,
                           y_max=None, fmt="0.00", note=None):
    data = CategoryChartData()
    data.categories = categories
    data.add_series("value", values)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, data)
    chart = gframe.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.gap_width = 60
    series = plot.series[0]
    series.data_labels.number_format = fmt
    series.data_labels.number_format_is_linked = False
    series.data_labels.show_value = True
    series.data_labels.font.size = Pt(12)
    series.data_labels.font.bold = True
    series.data_labels.font.color.rgb = INK
    for i, pt in enumerate(series.points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colors[i]
        pt.format.line.fill.background()
    style_chart_axes(chart, y_max=y_max, y_fmt=fmt)
    if note:
        _, tf = add_textbox(slide, x, Emu(y + h) + Inches(0.05), w, Inches(0.3))
        p = tf.paragraphs[0]
        add_inline_runs(p, note, 10.5, INK_MUTED, italic=True)
    return gframe


def add_grouped_bar(slide, x, y, w, h, categories, series_dict, series_colors,
                     y_max=None, fmt="0.00", note=None):
    data = CategoryChartData()
    data.categories = categories
    for name, vals in series_dict.items():
        data.add_series(name, vals)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, data)
    chart = gframe.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(12)
    chart.legend.font.color.rgb = INK_SECONDARY
    plot = chart.plots[0]
    plot.gap_width = 90
    plot.overlap = -10
    for i, series in enumerate(plot.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = series_colors[i]
        series.format.line.fill.background()
        series.data_labels.number_format = fmt
        series.data_labels.number_format_is_linked = False
        series.data_labels.show_value = True
        series.data_labels.font.size = Pt(10.5)
        series.data_labels.font.bold = True
        series.data_labels.font.color.rgb = INK
    style_chart_axes(chart, y_max=y_max, y_fmt=fmt)
    if note:
        _, tf = add_textbox(slide, x, Emu(y + h) + Inches(0.05), w, Inches(0.3))
        p = tf.paragraphs[0]
        add_inline_runs(p, note, 10.5, INK_MUTED, italic=True)
    return gframe


def add_horiz_bar(slide, x, y, w, h, categories, values, colors, fmt="0%", note=None):
    data = CategoryChartData()
    data.categories = categories
    data.add_series("value", values)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, w, h, data)
    chart = gframe.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.gap_width = 50
    series = plot.series[0]
    series.data_labels.number_format = fmt
    series.data_labels.number_format_is_linked = False
    series.data_labels.show_value = True
    series.data_labels.font.size = Pt(12)
    series.data_labels.font.bold = True
    series.data_labels.font.color.rgb = INK
    for i, pt in enumerate(series.points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colors[i]
        pt.format.line.fill.background()
    cat_ax = chart.category_axis
    cat_ax.tick_labels.font.size = Pt(13)
    cat_ax.tick_labels.font.color.rgb = INK
    cat_ax.format.line.color.rgb = GRIDLINE
    val_ax = chart.value_axis
    val_ax.visible = False
    if note:
        _, tf = add_textbox(slide, x, Emu(y + h) + Inches(0.05), w, Inches(0.3))
        p = tf.paragraphs[0]
        add_inline_runs(p, note, 10.5, INK_MUTED, italic=True)
    return gframe


def connector(slide, x1, y1, x2, y2, color=INK_MUTED, width=Pt(1.75)):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = width
    ln = conn.line._get_or_add_ln()
    from pptx.oxml.ns import qn as _qn
    tail = ln.makeelement(_qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    return conn


def flow_box(slide, x, y, w, h, label, sublabel=None, fill=SURFACE, line=BLUE,
             text_color=INK, size=13):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1.5)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_inline_runs(p, label, size, text_color)
    for r in p.runs:
        r.font.bold = True
    if sublabel:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        add_inline_runs(p2, sublabel, size - 3, INK_MUTED)
    return shp


# ======================================================================
# BUILD DECK
# ======================================================================

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH

SLIDE_NOTES = []  # (slide, text) filled inline as we go


# ---------------------------------------------------------------- 1. TITLE
s = new_slide(prs, bg=WHITE)
add_rect(s, 0, 0, SW, Inches(0.14), fill=BLUE)
add_rect(s, 0, SH - Inches(0.14), SW, Inches(0.14), fill=BLUE)
_, tf = add_textbox(s, Inches(1.0), Inches(2.35), Inches(11.3), Inches(0.4))
p = tf.paragraphs[0]
add_inline_runs(p, "FINAL POST-GRADUATION PROJECT · DEFENSE PRESENTATION", 13, BLUE)
p.runs[0].font.bold = True

_, tf = add_textbox(s, Inches(1.0), Inches(2.85), Inches(11.3), Inches(1.6))
p = tf.paragraphs[0]
add_inline_runs(p, "MemoryLifeBench", 44, INK)
p.runs[0].font.bold = True
p2 = tf.add_paragraph()
add_inline_runs(p2, "Memory Lifetime Prediction as a Time-to-Event Problem", 22, INK_SECONDARY)
p2.space_before = Pt(6)

add_rect(s, Inches(1.0), Inches(4.55), Inches(1.4), Pt(3), fill=BLUE)

_, tf = add_textbox(s, Inches(1.0), Inches(4.8), Inches(9), Inches(0.9))
p = tf.paragraphs[0]
add_inline_runs(p, "Bhargav Shendge", 18, INK)
p.runs[0].font.bold = True
p2 = tf.add_paragraph()
add_inline_runs(p2, "6-Week Independent Research Project", 14, INK_MUTED)
p2.space_before = Pt(2)

set_notes(s, "Thank you. Today I'm presenting my final post-graduation project: reframing memory-lifetime "
              "prediction for AI assistants as a survival-analysis problem, not a heuristic importance score. "
              "I'll walk through the full arc -- a model that wins on its own metric, a real deployment test that "
              "shows that isn't enough, the root-cause diagnosis, the fix, and two honesty checks that kept the "
              "final conclusion from overselling itself.")

# ---------------------------------------------------------------- 2. AGENDA
s = new_slide(prs)
kicker_and_title(s, "Overview", "Agenda")
agenda = [
    "Motivation & the gap in existing memory systems",
    "MemoryLifeBench: dataset and problem framing",
    "Method: the survival model and its joint multi-task extension",
    "Part 1 results: does it rank memory lifetimes well?",
    "Part 2: does that ranking translate into a good deployed policy?",
    "The surprise negative result → root-cause diagnosis → the fix",
    "Two honesty checks: an Oracle ceiling and a real Mem0 comparison",
    "Limitations, conclusion, and questions",
]
items = [{"text": f"**{i+1}.**  {t}", "size": 18} for i, t in enumerate(agenda)]
add_bullets(s, items, Inches(0.8), Inches(1.75), Inches(11.5), Inches(5), size=18, space_after=16)
set_notes(s, "Quick roadmap: I'll move through the problem, the dataset, the method, then the two-part results "
              "story -- ranking quality first, then whether that ranking quality actually helps -- and close with "
              "limitations and two checks that kept the conclusions honest.")

# ---------------------------------------------------------------- 3. MOTIVATION
s = new_slide(prs)
kicker_and_title(s, "1 · Motivation", "Why memory lifetime, and why now?")
left = [
    {"text": "Long-running AI assistants accumulate facts about the user across every session", "size": 17},
    {"text": "Storage isn't free, and unbounded context windows don't solve it", "size": 17},
    {"text": "Retrieval quality degrades as the memory store grows", "size": 17},
    {"text": "Stale or contradicted facts actively hurt answers if never removed", "size": 17},
]
add_bullets(s, left, Inches(0.8), Inches(1.7), Inches(6.1), Inches(4.5), space_after=18)

box = add_rect(s, Inches(7.3), Inches(1.85), Inches(5.2), Inches(3.9), fill=SURFACE, line=GRIDLINE, line_w=Pt(1))
_, tf = add_textbox(s, Inches(7.6), Inches(2.05), Inches(4.6), Inches(3.5))
p = tf.paragraphs[0]
add_inline_runs(p, "The core question", 15, BLUE)
p.runs[0].font.bold = True
p2 = tf.add_paragraph()
add_inline_runs(p2, "“Should this memory be forgotten?”", 16, INK)
p2.space_before = Pt(10)
p3 = tf.add_paragraph()
add_inline_runs(p3, "is the wrong question.", 16, INK_MUTED, italic=True)
p3.space_before = Pt(2)
p4 = tf.add_paragraph()
add_inline_runs(p4, "“When will this memory stop being useful?”", 16, INK)
p4.space_before = Pt(14)
for r in p4.runs:
    r.font.bold = True
p5 = tf.add_paragraph()
add_inline_runs(p5, "is the question this project answers directly.", 16, INK_MUTED)
p5.space_before = Pt(2)

set_notes(s, "Personal AI assistants keep extracting facts about you across sessions -- where you live, what "
              "you're working on. That store keeps growing, and you can't just keep everything: retrieval gets "
              "noisier, and stale facts don't just sit there, they actively produce wrong answers when they "
              "contradict something newer. At some point the system has to decide what to forget and when -- and "
              "existing systems answer a slightly different, weaker question than the one that actually matters.")

# ---------------------------------------------------------------- 4. RELATED WORK GAP
s = new_slide(prs)
kicker_and_title(s, "1 · Related Work", "How existing systems answer this")
headers = ["System", "Approach", "Predicts an explicit time-to-event?"]
rows = [
    ["Generative Agents (Park et al., 2023)", "Hand-tuned recency/importance/relevance score, decayed", "No"],
    ["MemGPT (Packer et al., 2023)", "OS-style paging between working set and external storage", "No"],
    ["Mem0 (Chhikara et al., 2025)", "Heuristic consolidation; frequently discards by internal rule", "No"],
    ["MemoryBank (Zhong et al., 2024)", "Fixed time-decay curve for consolidation/forgetting", "No"],
    ["HippoRAG / GraphRAG / Zep / A-Mem", "Graph-structured retrieval; no temporal-per-memory model", "No"],
    ["REMem (Shu et al., 2026, ICLR)", "Time-aware episodic graph + agentic retrieval", "Partially (representation, not a trained target)"],
]
add_table(s, Inches(0.55), Inches(1.65), Inches(12.2), Inches(4.3), headers, rows,
          col_widths=[3.1, 4.4, 3.0], font_size=13)
_, tf = add_textbox(s, Inches(0.55), Inches(6.15), Inches(12.2), Inches(0.5))
p = tf.paragraphs[0]
add_inline_runs(p, "**None of these systems predict an explicit, evaluable time-to-event target for individual memories — the gap this work targets.**", 14.5, INK)

set_notes(s, "Every existing memory system answers this with some flavor of scalar importance -- a recency-"
              "weighted heuristic, a paging scheme, a consolidation rule, or a graph structure. None of them "
              "predict an actual time: none say 'this fact will likely stop mattering in about three weeks.' "
              "REMem, the closest ICLR 2026 paper, is time-aware in its representation but doesn't train an "
              "evaluable lifetime target the way this project does.")

# ---------------------------------------------------------------- 5. REFRAMING
s = new_slide(prs)
kicker_and_title(s, "1 · Our Reframing", "Memory retention as survival analysis")
add_bullets(s, [
    {"text": "Treat every memory as a **survival-analysis subject**", "size": 18},
    {"text": "Born when stated → dies when invalidated, contradicted, superseded, or never referenced again", "size": 17, "level": 1},
    {"text": "Right-censored if the “death” is never observed in the data window", "size": 17, "level": 1},
    {"text": "Inherits a mature, well-understood toolkit for free", "size": 18},
    {"text": "Cox proportional hazards · concordance index · time-dependent AUC · integrated Brier score", "size": 17, "level": 1},
    {"text": "A sharp distinction the field draws that this paper leans on: **ranking quality (discrimination) is not the same as calibration on an absolute cutoff**", "size": 17, "level": 1},
], Inches(0.8), Inches(1.7), Inches(11.6), Inches(5), space_after=14)

set_notes(s, "We reframe this as survival analysis -- the same statistical framework used for patient survival "
              "times and equipment failure times. Every memory is born when stated, and dies when it's invalidated "
              "or never referenced again. Often we don't observe the death within our data window -- that's "
              "censoring, handled correctly rather than ignored. This buys us a mature toolkit for free. And it "
              "draws a sharp distinction I want you to hold onto: a model can rank very well and still be badly "
              "calibrated on an absolute cutoff. That distinction is the hinge the second half of this talk turns on.")

# ---------------------------------------------------------------- 6. CONTRIBUTIONS
s = new_slide(prs)
kicker_and_title(s, "1 · Contributions", "Four contributions")
contribs = [
    ("1", "MemoryLifeBench", "A 10,152-record benchmark of memory statements with time-to-event labels — synthetic dialogues with provable lifetimes plus real conversations from LoCoMo and LongMemEval, split by conversation with verified no leakage."),
    ("2", "A lightweight survival model", "213,889 parameters, frozen sentence embedding + a small MLP, that beats three frontier/local LLM baselines and two heuristics at ranking memory lifetimes."),
    ("3", "A joint multi-task extension", "Lifetime + Action + Future-Utility heads sharing one fused representation — still under 500K parameters — improving ranking further and adding two usable deployment signals."),
    ("4", "An honest downstream test", "A matched-storage-budget QA evaluation that catches a real failure the ranking metric could not, a root-cause diagnosis, two fixes, and two honesty checks against an Oracle ceiling and a real independent memory system."),
]
y = Inches(1.65)
for num, title, desc in contribs:
    add_rect(s, Inches(0.75), y, Inches(0.55), Inches(0.55), fill=BLUE)
    _, tf = add_textbox(s, Inches(0.75), y, Inches(0.55), Inches(0.55), anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_inline_runs(p, num, 20, WHITE)
    p.runs[0].font.bold = True
    _, tf2 = add_textbox(s, Inches(1.55), y - Inches(0.03), Inches(10.9), Inches(1.15))
    p2 = tf2.paragraphs[0]
    add_inline_runs(p2, title, 17, INK)
    p2.runs[0].font.bold = True
    p3 = tf2.add_paragraph()
    add_inline_runs(p3, desc, 13.5, INK_SECONDARY)
    p3.space_before = Pt(2)
    y = y + Inches(1.28)

set_notes(s, "Four contributions in total: the benchmark itself, the lightweight survival model, the joint "
              "multi-task extension, and -- the part I'd argue matters most -- an honest downstream test that "
              "catches a real failure, diagnoses it, fixes it, and then checks the fix against two more honesty "
              "tests rather than stopping at the first result that looked good.")

# ---------------------------------------------------------------- 7. DATASET
s = new_slide(prs)
kicker_and_title(s, "2 · Dataset", "MemoryLifeBench: composition")
headers = ["Source", "Train", "Val", "Test", "What it is"]
rows = [
    ["Synthetic", "3,199", "256", "265", "LLM-generated dialogues, exact provable lifetimes"],
    ["LongMemEval", "3,162", "359", "375", "Candidate memories from LongMemEval dialogues"],
    ["LoCoMo", "1,936", "324", "276", "Candidate memories from LoCoMo dialogues"],
    ["Total", "8,297", "939", "916", ""],
]
add_table(s, Inches(0.55), Inches(1.65), Inches(12.2), Inches(2.7), headers, rows,
          col_widths=[2.0, 1.3, 1.3, 1.3, 6.3], font_size=13.5, highlight_rows={3})
add_bullets(s, [
    {"text": "Split **by conversation** (568 / 71 / 71) — verified disjoint, no leakage across splits", "size": 15},
    {"text": "Event rate: 37.2% train / 35.5% val / 36.6% test", "size": 15},
    {"text": "Censoring convention is a disclosed judgment call for real conversations (Section 3.3) — flagged as a limitation, not presented as settled", "size": 15},
], Inches(0.55), Inches(4.7), Inches(12.2), Inches(2), space_after=10)

set_notes(s, "Just over 10,000 memory records with time-to-event labels, combining synthetic dialogues with a "
              "known true lifetime with real conversations from two published benchmarks. Everything is split by "
              "conversation, not by record, and verified to have no leakage. Worth flagging now, since it's an "
              "honest limitation: the censoring convention for real conversations is a judgment call, not a "
              "ground-truth fact -- I'll come back to that.")

# ---------------------------------------------------------------- 8. METHOD: SURVIVAL MODEL DIAGRAM
s = new_slide(prs)
kicker_and_title(s, "3 · Method", "Survival model (single-task)")
labels = [
    ("Memory\ntext", None),
    ("Frozen BGE\nembedding (768-d)", None),
    ("Small MLP\n(213,889 params)", None),
    ("Log hazard\nscore", None),
]
box_w, box_h, gap = Inches(2.55), Inches(1.1), Inches(0.55)
x0 = Inches(0.6)
y0 = Inches(2.3)
xs = []
for i, (lab, sub) in enumerate(labels):
    x = Emu(x0 + i * (box_w + gap))
    xs.append(x)
    fill = BLUE if i == len(labels) - 1 else SURFACE
    txt_color = WHITE if i == len(labels) - 1 else INK
    line_c = BLUE
    flow_box(s, x, y0, box_w, box_h, lab, fill=fill, line=line_c, text_color=txt_color, size=14)
for i in range(len(labels) - 1):
    connector(s, Emu(xs[i] + box_w), Emu(y0 + box_h // 2), xs[i + 1], Emu(y0 + box_h // 2), color=INK_MUTED, width=Pt(2))

add_bullets(s, [
    {"text": "No fine-tuning of the embedding, no generation, nothing autoregressive", "size": 16},
    {"text": "Trained with the standard **Cox partial-likelihood loss** — correctly handles right-censored durations", "size": 16},
    {"text": "**Zero token cost by construction** — a real, measurable advantage over every LLM-prompted baseline", "size": 16},
], Inches(0.6), Inches(4.1), Inches(11.8), Inches(2.5), space_after=14)

set_notes(s, "The core model is deliberately small. A frozen sentence embedding feeds a 214-thousand-parameter "
              "MLP that outputs a single hazard score, trained with the standard Cox partial-likelihood objective, "
              "which correctly handles the censored records rather than throwing them away. No fine-tuning, no "
              "generation -- which matters, because this model costs zero tokens to run, versus real metered API "
              "spend for the LLM baselines I'll show shortly.")

# ---------------------------------------------------------------- 9. METHOD: JOINT MODEL DIAGRAM
s = new_slide(prs)
kicker_and_title(s, "3 · Method", "Joint multi-task extension")
feat_labels = ["Intent", "Entities/NER", "Temporal", "Emotion/\nPreference", "Novelty", "Contradiction", "BGE\nembedding"]
fx = Inches(0.55)
fy = Inches(1.75)
fw, fh, fgap = Inches(1.55), Inches(0.62), Inches(0.14)
fusion_x = Inches(5.85)
fusion_y = Inches(2.85)
fusion_w, fusion_h = Inches(2.0), Inches(1.1)
for i, lab in enumerate(feat_labels):
    y = Emu(fy + i * (fh + fgap))
    flow_box(s, fx, y, fw, fh, lab, fill=SURFACE, line=AQUA, text_color=INK, size=11)
    connector(s, Emu(fx + fw), Emu(y + fh // 2), fusion_x, Emu(fusion_y + fusion_h // 2), color=RGBColor(0xC3, 0xC2, 0xB7), width=Pt(1))
flow_box(s, fusion_x, fusion_y, fusion_w, fusion_h, "Fusion", "(concat or gated)", fill=BLUE, line=BLUE, text_color=WHITE, size=15)

head_labels = [("Lifetime\nhead", "Cox survival objective", BLUE), ("Action\nhead", "4-way: store/update/\nmerge/forget", ORANGE), ("Future-Utility\nhead", "P(retrieved again)", MAGENTA)]
hx = Inches(9.1)
hw, hh, hgap = Inches(3.3), Inches(1.0), Inches(0.35)
hy0 = Inches(1.75)
for i, (lab, sub, color) in enumerate(head_labels):
    y = Emu(hy0 + i * (hh + hgap))
    flow_box(s, hx, y, hw, hh, lab, sub, fill=SURFACE, line=color, text_color=INK, size=14)
    connector(s, Emu(fusion_x + fusion_w), Emu(fusion_y + fusion_h // 2), hx, Emu(y + hh // 2), color=RGBColor(0xC3, 0xC2, 0xB7), width=Pt(1))

_, tf = add_textbox(s, Inches(0.55), Inches(6.15), Inches(12.2), Inches(0.9))
p = tf.paragraphs[0]
add_inline_runs(p, "Six frozen auxiliary extractors (~433M params, single forward pass, zero generation) fused with BGE → **three heads, one shared representation, 425,734 trainable params total**", 14, INK_SECONDARY)

set_notes(s, "The joint version fuses six off-the-shelf frozen feature extractors with the embedding, feeding "
              "three heads that share one representation: the same Lifetime head, an Action head classifying "
              "store, update, merge, or forget, and a Future-Utility head predicting whether a memory gets used "
              "again. All trained jointly with a custom loop, still under half a million trainable parameters "
              "total. Keep an eye on the Future-Utility head -- it becomes the hero of the second half of this talk.")

# ---------------------------------------------------------------- 10. METHOD: MEMORY SYSTEM PIPELINE
s = new_slide(prs)
kicker_and_title(s, "3 · Method", "The memory-lifecycle system")
stages = ["Extract &\nembed", "Store\n(vector DB)", "Forget\n(TTL + Action)", "Compact\n(merge dupes)", "Reflect\n(utility decay)", "Retrieve &\nrerank", "Grounded\nQA"]
n = len(stages)
box_w = Inches(1.62)
gap = Inches(0.2)
total_w = box_w * n + gap * (n - 1)
x0 = Emu(int((SW - total_w) / 2))
y0 = Inches(2.6)
box_h = Inches(1.1)
xs = []
for i, lab in enumerate(stages):
    x = Emu(x0 + i * (box_w + gap))
    xs.append(x)
    color = CAT8[i % len(CAT8)]
    flow_box(s, x, y0, box_w, box_h, lab, fill=SURFACE, line=color, text_color=INK, size=11.5)
for i in range(n - 1):
    connector(s, Emu(xs[i] + box_w), Emu(y0 + box_h // 2), xs[i + 1], Emu(y0 + box_h // 2), color=INK_MUTED, width=Pt(1.75))

add_bullets(s, [
    {"text": "Memory objects: {text, embedding, importance, predicted_ttl_days, action, utility_prob, provenance}", "size": 15},
    {"text": "Brute-force numpy vector store — sufficient at MemoryLifeBench's ~10K-memory scale", "size": 15},
    {"text": "Append-only audit log throughout; retriever reranks similarity-matched candidates by importance and utility before the grounded-QA call", "size": 15},
], Inches(0.6), Inches(4.35), Inches(11.8), Inches(2.2), space_after=12)

set_notes(s, "The trained heads feed a real memory-lifecycle system, not just an offline evaluation: extraction, "
              "storage in a vector store, a forgetting step driven by TTL and Action-head predictions, "
              "compaction, a reflection step that decays utility past predicted TTL, and a retriever that reranks "
              "before the final grounded-QA call. Everything is logged, append-only, for auditability.")

# ---------------------------------------------------------------- 11. SETUP PART 1
s = new_slide(prs)
kicker_and_title(s, "4 · Experiments, Part 1", "Setup: ranking quality")
add_bullets(s, [
    {"text": "Same test split (N=916) and validation split (N=939) for every method", "size": 17},
    {"text": "Five baselines:", "size": 17},
    {"text": "Recency-frequency heuristic", "size": 16, "level": 1},
    {"text": "Day/week/permanent bucket classifier", "size": 16, "level": 1},
    {"text": "Three LLM-prompted-TTL baselines: local Qwen2.5-7B, GPT-4o, Gemini 2.5 Pro (identical prompt across all three — “which model” is the only variable)", "size": 16, "level": 1},
    {"text": "Real, metered API spend for the two frontier baselines: 258,785 tokens (GPT-4o) + 210,072 tokens (Gemini)", "size": 17},
    {"text": "Our pipeline: **zero token cost by construction**", "size": 17, "bold": True},
], Inches(0.8), Inches(1.7), Inches(11.6), Inches(5), space_after=12)

set_notes(s, "Five baselines, all on the same held-out splits: a recency-frequency heuristic, a bucket "
              "classifier, and three LLM-prompted baselines using the exact same prompt so the model is the only "
              "variable being compared. Real, metered token spend for the two frontier baselines is reported "
              "directly. Our own pipeline costs zero tokens by construction.")

# ---------------------------------------------------------------- 12. HEADLINE C-INDEX CHART
s = new_slide(prs)
kicker_and_title(s, "4 · Experiments, Part 1", "Headline result: ranking quality (C-index)")
cats = ["Our model", "Bucket\nclassifier", "GPT-4o", "Qwen2.5-7B", "Gemini\n2.5 Pro", "Recency\nheuristic"]
vals = [0.7218, 0.6298, 0.5411, 0.5207, 0.4806, 0.4753]
colors = [BASELINE_COLOR[c.replace("\n", " ")] for c in cats]
add_single_series_bar(s, Inches(0.9), Inches(1.75), Inches(11.5), Inches(4.1), cats, vals, colors, y_max=0.8,
                       fmt="0.00")
_, tf = add_textbox(s, Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.7))
p = tf.paragraphs[0]
add_inline_runs(p, "**Beats every baseline, p < 0.001 (bootstrap, all pairwise).** vs. GPT-4o: +0.180 C-index, 95% CI [+0.134, +0.224]", 14.5, INK)

set_notes(s, "Here's the headline result of Part 1. Our model beats every baseline -- including prompting GPT-4o "
              "and Gemini 2.5 Pro directly -- by a wide, statistically significant margin. Every comparison clears "
              "p less than 0.001. Against GPT-4o alone that's an 18-point C-index gap, and it costs zero tokens.")

# ---------------------------------------------------------------- 13. RICHER METRICS + DETERMINISM
s = new_slide(prs)
kicker_and_title(s, "4 · Experiments, Part 1", "Richer metrics & determinism")
headers = ["Metric", "Our model", "GPT-4o"]
rows = [
    ["Time-dependent AUC", "0.7895", "0.5339"],
    ["Integrated Brier score (lower is better)", "0.2463", "0.4882"],
    ["Coefficient of variation across paraphrases", "0.0000", "0.3845 (44% of records CV>0.5)"],
]
add_table(s, Inches(0.7), Inches(1.75), Inches(11.9), Inches(2.4), headers, rows,
          col_widths=[4.5, 3.0, 4.4], font_size=14)
_, tf = add_textbox(s, Inches(0.7), Inches(4.5), Inches(11.9), Inches(2.2))
p = tf.paragraphs[0]
add_inline_runs(p, "**Determinism matters independently of accuracy.** Reword the same fact three ways: our model gives the identical answer every time. Qwen2.5-7B's CV is 0.7456 — 85% of records vary substantially on semantically identical inputs.", 16, INK)
p2 = tf.add_paragraph()
add_inline_runs(p2, "A deployed forgetting policy that changes its mind on a reworded version of the same fact is a liability regardless of its average accuracy.", 15, INK_SECONDARY, italic=True)
p2.space_before = Pt(10)

set_notes(s, "Beyond the headline C-index, richer metrics agree: better time-dependent AUC, lower Brier score. "
              "And a second, orthogonal argument: paraphrase determinism. Reword the same fact three ways and our "
              "model gives the identical answer every time. GPT-4o and especially the local Qwen model vary "
              "substantially on semantically identical inputs -- a real liability for a deployed system.")

# ---------------------------------------------------------------- 14. STABILITY / ABLATIONS / JOINT MODEL
s = new_slide(prs)
kicker_and_title(s, "4 · Experiments, Part 1", "Stability, ablations, and the joint model")
cats = ["Single-task\n(5 seeds)", "+ BGE-large\nencoder", "Joint model,\ngated fusion", "Joint model,\nconcat fusion"]
vals = [0.7312, 0.7382, 0.7304, 0.7553]
errs = [0.0131, 0.0048, 0.0082, 0.0045]
colors = [BLUE, AQUA, YELLOW, MAGENTA]
add_single_series_bar(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(3.6), cats, vals, colors, y_max=0.85,
                       fmt="0.00")
add_bullets(s, [
    {"text": "Test C-index ± std across 5 (single-task) / 3 (joint) seeds — low variance, not a lucky draw", "size": 14},
    {"text": "Concat fusion (0.7553 ± 0.0045) beats gated fusion (0.7304 ± 0.0082) despite gated being more parameter-rich — noted, not over-explained", "size": 14},
], Inches(0.9), Inches(5.65), Inches(11.5), Inches(1.5), space_after=8)

set_notes(s, "This isn't a lucky single seed -- five seeds land in a tight band. Swapping in a larger encoder "
              "gives a small, within-noise difference. The joint multi-task model improves further, to 0.7553, "
              "with an even tighter spread. Interestingly, simple concatenation beats a learned gating mechanism "
              "-- we note that rather than over-explain it, since the gap is plausible overfitting at this "
              "dataset's scale, not a settled finding.")

# ---------------------------------------------------------------- 15. ACTION & UTILITY HEADS
s = new_slide(prs)
kicker_and_title(s, "4 · Experiments, Part 1", "The Action and Future-Utility heads")
headers = ["Head", "Aggregate", "Detail"]
rows = [
    ["Action (4-way)", "86.4% accuracy", "Minority classes (update/merge/forget): precision 0.41–0.69, recall = 1.000 on all three across all 6 checkpoints"],
    ["Future-Utility (binary)", "AUC 0.71–0.77", "Genuinely predictive, well above the 0.5 random baseline; noisier than the other two heads"],
]
add_table(s, Inches(0.55), Inches(1.75), Inches(12.2), Inches(2.3), headers, rows,
          col_widths=[2.4, 2.2, 7.6], font_size=13.5)
_, tf = add_textbox(s, Inches(0.55), Inches(4.4), Inches(12.2), Inches(2.2))
p = tf.paragraphs[0]
add_inline_runs(p, "Inverse-frequency class weighting means the Action head **never misses** a true forget/update/merge case, at the cost of over-flagging some “store” records — a defensible operating point (a missed forget is worse than an extra review flag).", 15.5, INK)
p2 = tf.add_paragraph()
add_inline_runs(p2, "This precision gap on the forget class becomes directly relevant to the downstream failure diagnosed next.", 15.5, INK)
p2.space_before = Pt(10)
for r in p2.runs:
    r.font.bold = True
    r.font.color.rgb = BLUE

set_notes(s, "Two more heads worth a beat. The Action head hits 86.4% aggregate accuracy, but recall is exactly "
              "1.0 on the rare update, merge, and forget classes across every checkpoint -- it never misses a true "
              "case, at the cost of some false positives. That precision gap becomes directly relevant in a "
              "moment. The Future-Utility head is genuinely predictive too, though noisier -- it's about to become "
              "the most important signal in this whole talk.")

# ---------------------------------------------------------------- 16. BRIDGE SLIDE
s = new_slide(prs, bg=BLUE)
_, tf = add_textbox(s, Inches(1.0), Inches(2.7), Inches(11.3), Inches(2.0), anchor=MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]
add_inline_runs(p, "Good ranking. But does that make a good policy?", 34, WHITE)
p.runs[0].font.bold = True
p2 = tf.add_paragraph()
add_inline_runs(p2, "A high C-index proves the model ranks memories well. It does not prove that the resulting deployed policy — what actually gets kept or deleted — preserves downstream answer quality.", 17, RGBColor(0xE3, 0xEE, 0xFB))
p2.space_before = Pt(20)

set_notes(s, "So far, everything has been validated on the model's own ranking metric. But a memory system "
              "doesn't rank memories -- it deletes some of them. Does good ranking actually translate into a good "
              "deployed policy? That's the harder, and much more honest, question the rest of this talk answers.")

# ---------------------------------------------------------------- 17. SETUP PART 2
s = new_slide(prs)
kicker_and_title(s, "5 · Experiments, Part 2", "Setup: the matched-budget downstream test")
add_bullets(s, [
    {"text": "Four policies compared at the **same final storage budget per conversation**:", "size": 17},
    {"text": "`no_forget` — retain everything (a ceiling reference)", "size": 16, "level": 1},
    {"text": "`fifo` — keep the N most-recently-created memories", "size": 16, "level": 1},
    {"text": "`lru` — keep the N most-recently-referenced memories", "size": 16, "level": 1},
    {"text": "`ours` — our original policy: Lifetime-head TTL expiry + Action-head “forget”", "size": 16, "level": 1},
    {"text": "fifo/lru's capacity N is set to `ours`'s own natural budget — apples-to-apples, not an arbitrary cap", "size": 17},
    {"text": "Real LoCoMo + LongMemEval questions, answered by GPT-4o over the retrieved store, scored by exact match (EM) and token-F1", "size": 17},
], Inches(0.8), Inches(1.7), Inches(11.6), Inches(5), space_after=12)

set_notes(s, "To test this honestly, every policy is matched to the same final storage budget per conversation. "
              "fifo and lru get exactly the same capacity our own policy naturally settles at -- so this is a fair, "
              "apples-to-apples comparison, not an arbitrary fixed budget. Real questions, answered by GPT-4o over "
              "whatever memories survived, scored by exact match and F1.")

# ---------------------------------------------------------------- 18. SURPRISE NEGATIVE RESULT
s = new_slide(prs)
kicker_and_title(s, "5 · Experiments, Part 2", "The surprise: our policy was the worst")
cats = ["no_forget\n(ceiling)", "fifo", "lru", "ours"]
em = [0.1706, 0.1316, 0.1245, 0.1102]
f1 = [0.3007, 0.2347, 0.2319, 0.2054]
colors = [BLUE, ORANGE]
add_grouped_bar(s, Inches(1.0), Inches(1.75), Inches(11.3), Inches(4.0), cats,
                {"Mean EM": em, "Mean F1": f1}, colors, y_max=0.35, fmt="0.00")
_, tf = add_textbox(s, Inches(1.0), Inches(6.05), Inches(11.3), Inches(0.8))
p = tf.paragraphs[0]
add_inline_runs(p, "LoCoMo, N=1,542. **`ours` — built on the model that just won every ranking comparison — was the worst policy tested**, below both naive baselines and well below the no-forgetting ceiling.", 15, CRITICAL)
p.runs[0].font.bold = True

set_notes(s, "This is the moment the talk turns. On real downstream questions, our original policy -- the one "
              "built on the model that just won every ranking comparison -- was the worst policy tested. Worse "
              "than first-in-first-out. Worse than least-recently-used. That is the opposite of what the C-index "
              "results predicted, and we chased it down rather than smoothing it over.")

# ---------------------------------------------------------------- 19. ROOT CAUSE DIAGNOSIS
s = new_slide(prs)
kicker_and_title(s, "5 · Experiments, Part 2", "Root-cause diagnosis")
cats = ["no_forget", "lru", "fifo", "ours"]
vals = [1.0000, 0.7676, 0.7247, 0.6687]
colors = [POLICY_COLOR[c] for c in ["no_forget", "lru", "fifo", "ours"]]
add_single_series_bar(s, Inches(0.6), Inches(1.7), Inches(5.6), Inches(3.6), cats, vals, colors, y_max=1.15,
                       fmt="0.00", note="Evidence retention rate — LoCoMo, N=1,304 covered QA pairs")

_, tf = add_textbox(s, Inches(6.7), Inches(1.7), Inches(5.9), Inches(4.5))
p = tf.paragraphs[0]
add_inline_runs(p, "**Cause 1 — miscalibrated threshold**", 16, INK)
p.runs[0].font.bold = True
p2 = tf.add_paragraph()
add_inline_runs(p2, "TTL cutoff = survival curve's median = a coin flip by construction. Mean shortfall: 114.4 predicted days vs. 164.1 actually needed.", 14.5, INK_SECONDARY)
p2.space_before = Pt(4)
p3 = tf.add_paragraph()
add_inline_runs(p3, "Cause 2 — wrong decision structure", 16, INK)
p3.space_before = Pt(20)
for r in p3.runs:
    r.font.bold = True
p4 = tf.add_paragraph()
add_inline_runs(p4, "100% of losses came from TTL expiry, 0% from the Action head. An independent per-memory threshold, not a ranked top-N like fifo/lru use.", 14.5, INK_SECONDARY)
p4.space_before = Pt(4)

set_notes(s, "We built a free diagnostic -- no LLM calls -- checking whether the gold-evidence memory literally "
              "survives eviction. Our policy kept only two-thirds of what it needed. Two causes, both verified "
              "directly. One: the TTL cutoff was the survival curve's median -- a coin-flip cutoff by definition. "
              "Two, the deeper one: our policy makes an independent per-memory decision, while fifo and lru always "
              "keep their best N by a ranked score.")

# ---------------------------------------------------------------- 20. FIX 1: QUANTILE SWEEP
s = new_slide(prs)
kicker_and_title(s, "5 · Experiments, Part 2", "Fix #1: a configurable quantile cutoff")
headers = ["Q (S(t) cutoff)", "`ours` storage kept", "`ours` retention", "`lru` retention (same row)"]
rows = [
    ["0.5 (original)", "70.9%", "0.6687", "0.7676"],
    ["0.2", "91.3%", "0.9080", "0.9517"],
    ["0.1", "95.2%", "0.9502", "0.9709"],
    ["0.05", "97.1%", "0.9716", "—"],
]
add_table(s, Inches(0.8), Inches(1.85), Inches(11.5), Inches(2.9), headers, rows,
          col_widths=[2.5, 3.0, 3.0, 3.5], font_size=15)
_, tf = add_textbox(s, Inches(0.8), Inches(5.15), Inches(11.5), Inches(1.6))
p = tf.paragraphs[0]
add_inline_runs(p, "Making the cutoff a configurable quantile confirms cause 1 is real and fixable — but **does not close the gap against `lru` at any quantile tested**. `lru` retains evidence at least as well as `ours` at every row.", 16, INK)
p2 = tf.add_paragraph()
add_inline_runs(p2, "Cause 1 alone isn't the whole story — cause 2, the decision structure, is the residual gap.", 15, INK_SECONDARY, italic=True)
p2.space_before = Pt(8)

set_notes(s, "Making the TTL cutoff a configurable quantile instead of a hardcoded median lets us trade storage "
              "for retention, and the sweep confirms the fix direction. But it does not close the gap against lru "
              "at any quantile we tested -- lru retains evidence at least as well as our policy at every single "
              "row. That confirms cause one is real, but it leaves cause two, the decision structure, unresolved.")

# ---------------------------------------------------------------- 21. FIX 2: RANKED UTILITY EVICTION
s = new_slide(prs)
kicker_and_title(s, "5 · Experiments, Part 2", "Fix #2: rank by the Future-Utility head")
cats = ["0.5\n(original)", "0.2", "0.1"]
fifo_v = [0.7247, 0.9225, 0.9670]
lru_v = [0.7676, 0.9517, 0.9709]
ours_v = [0.6687, 0.9080, 0.9502]
util_v = [0.8765, 0.9663, 0.9877]
add_grouped_bar(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(3.9), cats,
                {"fifo": fifo_v, "lru": lru_v, "ours (original)": ours_v, "ours_utility (fix)": util_v},
                [ORANGE, AQUA, YELLOW, MAGENTA], y_max=1.15, fmt="0.00")
_, tf = add_textbox(s, Inches(0.7), Inches(5.85), Inches(11.9), Inches(1.0))
p = tf.paragraphs[0]
add_inline_runs(p, "**`ours_utility` — rank all memories by the already-trained Future-Utility score, keep the top-N — beats every baseline at every quantile tested, including the original unfixed setting.**", 15, GOOD)
p.runs[0].font.bold = True

set_notes(s, "The fix wasn't a bigger model -- it was a better use of a model we already had. We already had a "
              "Future-Utility head, previously only used to rerank retrieval, never consulted by eviction. "
              "Switching eviction to rank by that score and keep the top N -- the same structure fifo and lru "
              "already use -- beats every baseline at every storage budget tested, including the original "
              "unfixed setting.")

# ---------------------------------------------------------------- 22. WHY THE FIX WORKS
s = new_slide(prs)
kicker_and_title(s, "5 · Experiments, Part 2", "Why does it work? A mechanistic explanation")
cats = ["utility_prob\n(Future-Utility head)", "predicted_ttl_days\n(Lifetime head, median)"]
vals = [0.6709, 0.2852]
colors = [MAGENTA, YELLOW]
add_single_series_bar(s, Inches(1.5), Inches(1.75), Inches(9.3), Inches(3.6), cats, vals, colors, y_max=1.0,
                       fmt="0.00", note="Pooled AUC predicting real QA-evidence relevance, across all 2,536 LoCoMo memories. 0.50 = random chance.")
_, tf = add_textbox(s, Inches(1.5), Inches(5.85), Inches(9.3), Inches(1.2))
p = tf.paragraphs[0]
add_inline_runs(p, "`predicted_ttl_days` scores **below 0.5 — inversely correlated** with evidence relevance. The fix works not because ranking beats thresholding in general, but because it replaces a signal that is actively backwards with one that is genuinely predictive.", 15, INK)

set_notes(s, "We didn't stop at 'it works' -- we checked why. utility_prob has real predictive power for whether "
              "a memory is actual QA evidence, AUC 0.67. The old TTL-based signal scores 0.29 -- below chance, "
              "literally pointing the wrong direction. That's a complete causal explanation, not a correlation "
              "we're hoping holds.")

# ---------------------------------------------------------------- 23. CONFIRMATION WITH SIGNIFICANCE
s = new_slide(prs)
kicker_and_title(s, "5 · Experiments, Part 2", "Confirming the fix, honestly")
headers = ["Comparison (n=120)", "EM diff, 95% CI", "p (fix beats other)"]
rows = [
    ["ours_utility vs. fifo", "+0.0499 [+0.0167, +0.0917]", "0.002 — significant"],
    ["ours_utility vs. ours (original)", "+0.0415 [+0.0083, +0.0833]", "0.006 — significant"],
    ["ours_utility vs. no_forget", "+0.0000 [+0.0000, +0.0000]", "1.000 — tied"],
    ["ours_utility vs. lru", "+0.0083 [+0.0000, +0.0250]", "0.367 — not significant"],
]
add_table(s, Inches(0.55), Inches(1.75), Inches(12.2), Inches(2.9), headers, rows,
          col_widths=[4.3, 4.3, 3.6], font_size=14, highlight_rows={3})
_, tf = add_textbox(s, Inches(0.55), Inches(4.95), Inches(12.2), Inches(1.6))
p = tf.paragraphs[0]
add_inline_runs(p, "Significantly beats fifo and the original policy; statistically **tied** with the no-forget ceiling; **not** significantly better than lru at this sample size — reported plainly, even though the free-diagnostic sweep and the AUC analysis both point the same direction.", 15, INK)
p2 = tf.add_paragraph()
add_inline_runs(p2, "Confirmed on a 4th metric (BLEU-1) and via refusal-precision/F1, following REMem's own methodology.", 14, INK_SECONDARY, italic=True)
p2.space_before = Pt(8)

set_notes(s, "On a real sample of 120 GPT-4o-scored questions, with proper paired bootstrap testing, the fixed "
              "policy significantly beats fifo and the original policy, and is statistically tied with the "
              "no-forget ceiling. It is not significantly better than lru at this sample size, and we say that "
              "plainly rather than dress it up. It's also confirmed on a fourth metric family, BLEU-1, and via "
              "refusal precision, following REMem's own methodology for exactly this purpose.")

# ---------------------------------------------------------------- 24. ERROR TAXONOMY
s = new_slide(prs)
kicker_and_title(s, "5 · Experiments, Part 2", "What does the model actually get wrong?")
cats = ["False refusal", "Wrong value", "Incomplete answer", "Date off-by-one", "Reasoning error", "Likely judge error"]
vals = [0.54, 0.21, 0.15, 0.05, 0.02, 0.03]
colors = [CRITICAL, ORANGE, YELLOW, AQUA, VIOLET, INK_MUTED]
add_horiz_bar(s, Inches(1.4), Inches(1.75), Inches(9.8), Inches(3.9), cats, vals, colors, fmt="0%",
              note="100 sampled wrong (judge-negative) predictions, pooled across all 5 policies and both benchmarks — a single-pass manual read, not independently double-rated.")
_, tf = add_textbox(s, Inches(1.4), Inches(6.05), Inches(9.8), Inches(0.7))
p = tf.paragraphs[0]
add_inline_runs(p, "**False refusal dominates** — most of what looks like “the model got it wrong” is actually “the eviction policy removed the evidence, and the model correctly said so.”", 15, INK)
p.runs[0].font.bold = True

set_notes(s, "We sampled 100 wrong predictions across every policy and categorized each one, following REMem's "
              "own error-analysis methodology. False refusal is the dominant failure mode by a wide margin -- not "
              "wrong guesses, not date-arithmetic mistakes. Most of what looks like 'the model got it wrong' is, "
              "on inspection, 'the eviction policy removed the evidence, and the model correctly reported that.'")

# ---------------------------------------------------------------- 25. QUESTION-TYPE BREAKDOWN
s = new_slide(prs)
kicker_and_title(s, "5 · Experiments, Part 2", "Does the fix help every question type equally?")
cats = ["Single-Hop\n(N=54, judge score)", "Multi-Hop\n(N=48, EM)", "Temporal\n(N=16, EM)"]
fifo_v = [0.204, 0.021, 0.1875]
ours_v = [0.259, 0.021, 0.1875]
util_v = [0.537, 0.083, 0.1875]
add_grouped_bar(s, Inches(1.0), Inches(1.75), Inches(11.3), Inches(3.9), cats,
                {"fifo": fifo_v, "ours (original)": ours_v, "ours_utility (fix)": util_v},
                [ORANGE, YELLOW, MAGENTA], y_max=0.65, fmt="0.00")
_, tf = add_textbox(s, Inches(1.0), Inches(5.95), Inches(11.3), Inches(1.0))
p = tf.paragraphs[0]
add_inline_runs(p, "The effect is **concentrated, not uniform**: a clear win on Single-Hop and Multi-Hop; **zero measured effect on Temporal questions** (identical 0.1875 across all 5 policies). Sample sizes here are small — flagged, not hidden.", 14.5, INK)

set_notes(s, "The improvement isn't uniform across question types. It's a clear win on Single-Hop and Multi-Hop "
              "questions -- Multi-Hop actually matches the no-forget ceiling exactly. On Temporal questions, "
              "though, this sample shows zero measured effect at all, identical across every policy. We flag the "
              "small sample size directly rather than over-interpreting either result.")

# ---------------------------------------------------------------- 26. HONESTY CHECK 1: ORACLE CEILING
s = new_slide(prs)
kicker_and_title(s, "6 · Honesty Checks", "Is “no-forget” really the ceiling?")
cats = ["oracle\n(gold evidence only)", "full_context\n(uncapped retrieval)", "no_forget\n(for reference)"]
em = [0.1589, 0.1333, 0.0917]
f1 = [0.3473, 0.2980, 0.1957]
add_grouped_bar(s, Inches(1.0), Inches(1.7), Inches(11.3), Inches(3.9), cats,
                {"Mean EM": em, "Mean F1": f1}, [VIOLET, BLUE], y_max=0.45, fmt="0.00")
_, tf = add_textbox(s, Inches(1.0), Inches(5.9), Inches(11.3), Inches(1.1))
p = tf.paragraphs[0]
add_inline_runs(p, "LoCoMo, N=107 (matched). **Oracle significantly beats no_forget**: F1 +0.1416, p<0.0001. Retrieval quality — not eviction policy — is a real, independent bottleneck this paper does not fix.", 15, INK)

set_notes(s, "Is 'keep everything' really the best anything could do? No. An oracle given only the correct "
              "evidence, with no retrieval noise, significantly beats the no-forget ceiling. So ordinary top-5 "
              "retrieval over a large store is its own separate bottleneck -- one this paper's eviction fixes were "
              "never positioned to close.")

# ---------------------------------------------------------------- 27. HONESTY CHECK 2: MEM0
s = new_slide(prs)
kicker_and_title(s, "6 · Honesty Checks", "A real, independent memory system: Mem0")
cats = ["mem0", "no_forget", "ours_utility", "lru", "ours", "fifo"]
vals = [0.2275, 0.1957, 0.1949, 0.1998, 0.1575, 0.1294]
colors = [POLICY_COLOR[c] for c in cats]
add_single_series_bar(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(3.5), cats, vals, colors, y_max=0.28,
                       fmt="0.00", note="Mean F1, N=120 LoCoMo questions, real GPT-4o-scored answers.")
_, tf = add_textbox(s, Inches(0.9), Inches(5.55), Inches(11.5), Inches(1.5))
p = tf.paragraphs[0]
add_inline_runs(p, "Bootstrap significance: mem0 is **statistically tied** with no_forget, ours_utility, and lru (all p>0.13) — not proven superior. It significantly beats fifo (p<0.001) and the original ours policy (p=0.005).", 15, INK)
p2 = tf.add_paragraph()
add_inline_runs(p2, "Reported as measured — including that Mem0 used a weaker local extraction model, a budget-forced handicap in Mem0's favor to correct for, not against.", 14, INK_SECONDARY, italic=True)
p2.space_before = Pt(6)

set_notes(s, "We integrated Mem0, a real independently-built memory system, end to end -- real indexing, real "
              "extraction, not just a discussion in related work. Our best policy is statistically tied with it, "
              "not superior, and we report that as measured rather than reframing it to sound better. Worth "
              "noting: Mem0 was handicapped here by a weaker local extraction model, purely a budget constraint, "
              "so a fair fight could only make its showing stronger.")

# ---------------------------------------------------------------- 28. LIMITATIONS
s = new_slide(prs)
kicker_and_title(s, "7 · Limitations", "Stated directly, not deferred to an appendix")
lims = [
    "Mem0 comparison used a weaker local LLM for its own extraction, forced by budget — not a fairness advantage",
    "Retrieval quality, not just eviction policy, is a real, unaddressed bottleneck (the Oracle result)",
    "The censoring convention for real conversations is a judgment call with no ground-truth alternative",
    "`ours_utility` vs. `lru` is not statistically significant at n=120 real LLM-scored questions",
    "The Importance head is a documented heuristic, not learned — no ground-truth label exists",
    "No human validation of extracted labels has been performed; all labels are programmatically derived",
]
add_bullets(s, [{"text": t, "size": 16.5} for t in lims], Inches(0.8), Inches(1.75), Inches(11.6), Inches(4.8), space_after=16)

set_notes(s, "Being direct about what this doesn't yet show: the Mem0 comparison used a weaker local model, "
              "purely a budget constraint. Retrieval quality is a real bottleneck this paper doesn't fix. The "
              "censoring convention is a judgment call. The lru comparison needs a bigger sample to clear "
              "significance. The Importance head is a heuristic, not learned. And no human has yet validated the "
              "extracted labels.")

# ---------------------------------------------------------------- 29. CONCLUSION
s = new_slide(prs, bg=BLUE)
_, tf = add_textbox(s, Inches(1.0), Inches(1.0), Inches(11.3), Inches(0.6))
p = tf.paragraphs[0]
add_inline_runs(p, "CONCLUSION", 14, RGBColor(0xCD, 0xE2, 0xFB))
p.runs[0].font.bold = True
_, tf = add_textbox(s, Inches(1.0), Inches(1.5), Inches(11.3), Inches(4.8))
p = tf.paragraphs[0]
add_inline_runs(p, "Ranking quality alone did not guarantee a good deployed policy.", 27, WHITE)
p.runs[0].font.bold = True
p2 = tf.add_paragraph()
add_inline_runs(p2, "A downstream test caught what the ranking metric structurally could not. The fix that worked was not a bigger model — it was validating the actual deployed decision, not just the metric the model was trained on.", 17, RGBColor(0xE3, 0xEE, 0xFB))
p2.space_before = Pt(16)
p3 = tf.add_paragraph()
add_inline_runs(p3, "Two honesty checks — an Oracle ceiling and a real Mem0 comparison — kept this conclusion from overselling itself.", 17, RGBColor(0xE3, 0xEE, 0xFB))
p3.space_before = Pt(10)

set_notes(s, "To close: framing this as survival analysis gives a tiny, deterministic, zero-token model that "
              "beats prompting frontier LLMs at ranking memory lifetimes. But the real contribution is what "
              "happened next -- ranking quality did not automatically produce a good deployed policy, and the fix "
              "that worked was validating the actual decision a deployed system makes, not just the metric it was "
              "trained against. If there's one transferable lesson here, it's that one.")

# ---------------------------------------------------------------- 30. REFERENCES
s = new_slide(prs)
kicker_and_title(s, "References", "Key references", title_size=26)
refs = [
    "Chhikara, P. et al. (2025). Mem0: Building Production-Ready AI Agents with Scalable Long-Term Memory. arXiv:2504.19413.",
    "Cox, D. R. (1972). Regression Models and Life-Tables. JRSS Series B, 34(2), 187-220.",
    "Maharana, A. et al. (2024). Evaluating Very Long-Term Conversational Memory of LLM Agents (LoCoMo). arXiv:2402.17753.",
    "Packer, C. et al. (2023). MemGPT: Towards LLMs as Operating Systems. arXiv:2310.08560.",
    "Park, J. S. et al. (2023). Generative Agents: Interactive Simulacra of Human Behavior.",
    "Shu, Y. et al. (2026). REMem: Reasoning with Episodic Memory in Language Agents. ICLR 2026, arXiv:2602.13530.",
    "Wu, D. et al. (2024). LongMemEval: Benchmarking Chat Assistants on Long-Term Interactive Memory. arXiv:2410.10813.",
    "Full reference list (14 entries) in the accompanying manuscript.",
]
add_bullets(s, [{"text": t, "size": 13.5} for t in refs], Inches(0.8), Inches(1.75), Inches(11.6), Inches(4.8), space_after=12)

set_notes(s, "The full reference list of fourteen entries is in the accompanying manuscript; these are the six "
              "most load-bearing for this talk.")

# ---------------------------------------------------------------- 31. THANK YOU
s = new_slide(prs, bg=WHITE)
add_rect(s, 0, 0, SW, Inches(0.14), fill=BLUE)
add_rect(s, 0, SH - Inches(0.14), SW, Inches(0.14), fill=BLUE)
_, tf = add_textbox(s, Inches(1.0), Inches(2.9), Inches(11.3), Inches(1.2))
p = tf.paragraphs[0]
add_inline_runs(p, "Thank you.", 40, INK)
p.runs[0].font.bold = True
p2 = tf.add_paragraph()
add_inline_runs(p2, "Questions?", 22, INK_SECONDARY)
p2.space_before = Pt(10)
_, tf = add_textbox(s, Inches(1.0), Inches(4.6), Inches(11.3), Inches(0.6))
p3 = tf.paragraphs[0]
add_inline_runs(p3, "Code, data, and full paper: github.com/frustratedengineer-gif/pgp", 13, INK_MUTED)

set_notes(s, "Thank you -- happy to take questions. I have backup slides ready on the formal notation, the full "
              "policy definitions, and the extra significance tables if it's useful to go deeper on any of this.")

# ---------------------------------------------------------------- BACKUP SECTION DIVIDER
s = new_slide(prs, bg=INK)
_, tf = add_textbox(s, Inches(1.0), Inches(3.2), Inches(11.3), Inches(1.2), anchor=MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]
add_inline_runs(p, "Backup slides", 30, WHITE)
p.runs[0].font.bold = True
p2 = tf.add_paragraph()
add_inline_runs(p2, "For anticipated questions", 16, RGBColor(0xC3, 0xC2, 0xB7))
set_notes(s, "These backup slides are for anticipated follow-up questions -- formal notation, the full policy "
              "definitions, and extra tables not needed for the core narrative.")

# ---------------------------------------------------------------- B1: FORMAL NOTATION
s = new_slide(prs)
kicker_and_title(s, "Backup", "Formal notation", title_size=26)
add_bullets(s, [
    {"text": "Hazard:  h(t | z_i) = h_0(t) · exp(f(z_i))", "size": 16},
    {"text": "Survival:  S(t | z_i) = S_0(t) ^ exp(f(z_i))", "size": 16},
    {"text": "Cox partial log-likelihood over uncensored set D = {i : δ_i = 1}:", "size": 16},
    {"text": "L(θ) = Σᵢ [ f(z_i) − log( Σⱼ∈R(T_i) exp(f(z_j)) ) ],  R(t) = {j : T_j ≥ t}", "size": 15, "level": 1},
    {"text": "C-index over comparable pairs E = {(i,j): T_i<T_j, δ_i=1}: fraction correctly ordered", "size": 16},
    {"text": "Quantile TTL:  TTL_q(z_i) = sup{t : S(t | z_i) ≥ q}, capped at 3,650 days", "size": 16},
], Inches(0.8), Inches(1.75), Inches(11.6), Inches(4.5), space_after=14)

set_notes(s, "The formal Cox proportional-hazards notation, matching the actual implementation exactly, per "
              "Appendix A.1 of the manuscript.")

# ---------------------------------------------------------------- B2: ALL POLICIES
s = new_slide(prs)
kicker_and_title(s, "Backup", "All six eviction policies, formally", title_size=26)
headers = ["Policy", "Score / rule"]
rows = [
    ["no_forget", "Retain everything"],
    ["fifo", "top-N by created_at(o)"],
    ["lru", "top-N by last_referenced(o), −∞ if never referenced"],
    ["ours (original)", "action(o) ≠ forget AND age(o,t) ≤ predicted_ttl(o) — independent per-memory threshold"],
    ["ours_utility (the fix)", "top-N by utility_prob(o)"],
    ["ours_combo", "top-N by 0.5·utility_prob(o) + 0.5·remaining_life_fraction(o,t) — underperforms pure utility"],
]
add_table(s, Inches(0.55), Inches(1.75), Inches(12.2), Inches(4.0), headers, rows,
          col_widths=[3.0, 9.2], font_size=13.5)

set_notes(s, "The full formal definition of all six policies compared in this work, matching Appendix A.2 of the "
              "manuscript exactly.")

# ---------------------------------------------------------------- B3: MEM0 COST DETAIL
s = new_slide(prs)
kicker_and_title(s, "Backup", "Mem0 integration: cost & confounds", title_size=26)
add_bullets(s, [
    {"text": "Real calibration: 60 turns, $0.6381 spent → $0.01063/turn → **$62.55 extrapolated** for full GPT-4o indexing of all 5,882 turns", "size": 15.5},
    {"text": "Budget-infeasible for a self-funded project — substituted a free, self-hosted local Qwen2.5-7B for Mem0's own extraction calls only", "size": 15.5},
    {"text": "QA-answering step stayed on real GPT-4o throughout, for a fair comparison basis", "size": 15.5},
    {"text": "Two disclosed confounds:", "size": 15.5},
    {"text": "No working timestamp parameter in Mem0's open-source release — local model once resolved “yesterday” to 2026 instead of the conversation's true 2023", "size": 14.5, "level": 1},
    {"text": "195 of 5,882 indexing calls (3.3%) produced malformed JSON, silently contributing zero memories for that turn", "size": 14.5, "level": 1},
], Inches(0.8), Inches(1.75), Inches(11.6), Inches(5), space_after=10)

set_notes(s, "The full disclosure of how the Mem0 baseline was made affordable, and the two real, measured "
              "confounds that resulted -- both handicapping Mem0 relative to its normal deployment, never our own "
              "policies.")

# ---------------------------------------------------------------- B4: REPRODUCIBILITY
s = new_slide(prs)
kicker_and_title(s, "Backup", "Reproducibility", title_size=26)
add_bullets(s, [
    {"text": "5 seeds (survival model) / 3 seeds (joint model) — all reported means include standard deviation", "size": 17},
    {"text": "51 unit tests (31 added during the downstream-evaluation phase)", "size": 17},
    {"text": "Every prompt used in this work reproduced verbatim in the manuscript's Appendix C", "size": 17},
    {"text": "Code released under the MIT license; data under CC BY-NC 4.0 (attribution to LoCoMo and LongMemEval required)", "size": 17},
], Inches(0.8), Inches(1.9), Inches(11.6), Inches(4.2), space_after=18)

set_notes(s, "Everything needed to reproduce this work is committed: seeds, unit tests, verbatim prompts, and an "
              "open license on both code and data.")

# ---------------------------------------------------------------------- footers
slides = list(prs.slides.__iter__())
total = len(slides)
for i, sl in enumerate(slides, start=1):
    if i in (1, 16, 29, 32):  # title / bridge / conclusion / backup-divider slides: no footer
        continue
    add_footer(sl, i, total)

out_path = "paper/MemoryLifeBench_final_defense.pptx"
prs.save(out_path)
print(f"written -> {out_path}  ({total} slides)")
